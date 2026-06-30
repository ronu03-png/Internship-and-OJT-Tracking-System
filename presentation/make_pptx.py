"""
Generate a PowerPoint (.pptx) version of the InternTrack presentation.

Requirements:
    pip install python-pptx

Usage:
    python make_pptx.py

Output:
    presentation/interntrack_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
PRIMARY = RGBColor(37, 99, 235)       # #2563EB
PRIMARY_DARK = RGBColor(30, 58, 138)  # #1E3A8A
ACCENT = RGBColor(16, 185, 129)       # #10B981
DARK = RGBColor(15, 23, 42)            # #0F172A
TEXT = RGBColor(51, 65, 85)             # #334155
MUTED = RGBColor(100, 116, 139)         # #64748B
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_blank_slide():
    """Add a blank slide with the correct layout."""
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


def add_title_bar(slide, title_text, subtitle_text="", dark=False):
    """Add a top title bar to a slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_DARK if dark else WHITE
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else PRIMARY_DARK
    p.font.name = "Inter"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = WHITE if dark else MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(6)


def add_bullet_box(slide, left, top, width, height, bullets, title=None, font_size=16):
    """Add a text box with bullet points."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Inter"
        p.space_after = Pt(10)

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.font.name = "Inter"
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.3


def add_footer(slide, slide_num):
    """Add slide number and accent line."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(1), Inches(0.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num}"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.font.name = "Inter"


def add_cover_slide():
    slide = add_blank_slide()
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Subtitle badge
    badge = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.6))
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "Capstone Project Presentation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.1), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "InternTrack"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Web-Based Internship and OJT Placement & Tracking System"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.font.name = "Inter"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Powered by OJTConnect — helping students find work through direct contact with companies and city offices"
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.font.name = "Inter"
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(18)

    p4 = tf.add_paragraph()
    p4.text = "Presented to: HR Department, City Hall"
    p4.font.size = Pt(14)
    p4.font.color.rgb = WHITE
    p4.font.name = "Inter"
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(24)

    add_footer(slide, 1)


def add_title_defense_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Title Defense", "Proposed capstone title and focus")

    box = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Proposed Title:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = MUTED
    p.font.name = "Inter"

    p2 = tf.add_paragraph()
    p2.text = "InternTrack: A Web-Based Internship and On-the-Job Training (OJT) Placement & Tracking System"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = DARK
    p2.font.name = "Inter"
    p2.space_before = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "Subtitle: Powered by OJTConnect — helping students find work through direct contact with companies and city offices"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT
    p3.font.name = "Inter"
    p3.space_before = Pt(8)

    add_bullet_box(slide, 0.6, 3.5, 5.8, 3.2,
                   ["Connects colleges, students, supervisors, and city offices.",
                    "Replaces spreadsheets, paper forms, and scattered communication."],
                   title="What it is")

    add_bullet_box(slide, 6.9, 3.5, 5.8, 3.2,
                   ["Streamlines OJT placement and tracking.",
                    "Provides transparency and accountability for all stakeholders."],
                   title="Why it matters")

    add_footer(slide, 2)


def add_problem_statement_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Problem Statement")

    add_bullet_box(slide, 0.6, 1.9, 12.1, 3.6,
                   ["Students struggle to find verified companies and city offices that accept OJT placements.",
                    "Supervisors and coordinators track attendance, reports, and requirements on paper or spreadsheets.",
                    "Communication is fragmented across calls, emails, and text messages.",
                    "There is no centralized visibility of student progress and completion status.",
                    "Approval of documents, reports, and attendance takes too long."],
                   title="Current challenges in OJT management")

    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(239, 246, 255)
    highlight.line.color.rgb = PRIMARY

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Core problem: There is no unified, digital, department-based system that lets students find OJT openings, apply directly, and be monitored throughout their internship."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Inter"
    p.line_spacing = 1.2

    add_footer(slide, 3)


def add_objectives_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Objectives")

    add_bullet_box(slide, 0.6, 1.9, 5.8, 3.0,
                   ["Find and apply to OJT openings by department and course.",
                    "Submit requirements, log attendance, and submit reports.",
                    "Communicate with supervisors and view progress."],
                   title="For Students")

    add_bullet_box(slide, 6.9, 1.9, 5.8, 3.0,
                   ["Post openings and review student applications.",
                    "Verify attendance and evaluate performance.",
                    "Approve/reject reports and requirements."],
                   title="For Supervisors")

    add_bullet_box(slide, 0.6, 5.0, 5.8, 2.0,
                   ["Manage students, companies, and placements.",
                    "Monitor progress across departments.",
                    "Generate summary reports."],
                   title="For Coordinators")

    add_bullet_box(slide, 6.9, 5.0, 5.8, 2.0,
                   ["Manage users, roles, and settings.",
                    "Maintain audit logs and ensure data security."],
                   title="For Administrators")

    add_footer(slide, 4)


def add_scope_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Scope & Limitations")

    add_bullet_box(slide, 0.6, 1.9, 5.8, 4.6,
                   ["Four user roles: Admin, Coordinator, Supervisor, Student.",
                    "Department and course-based division.",
                    "OJT opening board, applications, and placement.",
                    "Attendance, reports, requirements, and messaging.",
                    "Dashboards, notifications, analytics, and audit trail.",
                    "Responsive design for desktop, tablet, and mobile."],
                   title="In Scope")

    add_bullet_box(slide, 6.9, 1.9, 5.8, 4.6,
                   ["QR code and selfie/GPS attendance are optional future features.",
                    "Email/SMS notifications are planned for future integration.",
                    "Dedicated mobile app is not included.",
                    "Offline functionality is not supported."],
                   title="Out of Scope / Limitations")

    add_footer(slide, 5)


def add_why_hr_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Why the HR Department of the City Hall?")

    reasons = [
        ("Central Hub", "HR is the entry point for all OJT applicants entering city government offices."),
        ("High Volume", "City Hall receives many OJT students from different colleges each semester."),
        ("Multi-Department", "Students can be deployed to treasury, health, engineering, social services, and more."),
        ("Supervision Need", "HR must track attendance, hours, and completion for endorsement."),
        ("Real Impact", "Digitalizing the process reduces paperwork and improves public service."),
        ("Feasibility", "HR has clear authority, workflows, and a need for an auditable system."),
    ]

    x_positions = [0.6, 4.55, 8.5]
    y_positions = [1.9, 4.4]
    for i, (title, desc) in enumerate(reasons):
        col = i % 3
        row = i // 3
        left = x_positions[col]
        top = y_positions[row]

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.7), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT
        p2.font.name = "Inter"
        p2.space_before = Pt(6)
        p2.line_spacing = 1.2

    add_footer(slide, 6)


def add_solution_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Proposed Solution")

    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.1)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(239, 246, 255)
    highlight.line.color.rgb = PRIMARY

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "OJTConnect is the deployment module of InternTrack. It gives students a searchable board of openings, lets them apply directly, and then tracks their progress once accepted."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Inter"
    p.line_spacing = 1.2

    steps = [
        ("1", "Discover", "Students browse openings filtered by department and course."),
        ("2", "Connect", "Students apply or contact the HR/supervisor directly via email."),
        ("3", "Track", "Attendance, reports, and requirements are tracked in one place."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        left = 0.6 + i * 4.2
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.3), Inches(3.9), Inches(2.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.45), Inches(3.5), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{num}. {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT
        p2.font.name = "Inter"
        p2.space_before = Pt(10)
        p2.line_spacing = 1.3

    add_footer(slide, 7)


def add_features_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Key Features")

    features = [
        ("Authentication", "Login, register, role-based access, and secure sessions."),
        ("Dashboards", "Role-specific views with stats, progress bars, and quick actions."),
        ("Attendance", "Time in/out logging, automatic hour computation, and approval."),
        ("Reports", "Narrative reports with supervisor feedback and approval."),
        ("Requirements", "Document submission and status tracking."),
        ("Messaging", "Built-in chat between students and supervisors."),
        ("Openings", "OJT posting board and application workflow."),
        ("Analytics", "Progress tracking, completion rates, and exportable reports."),
        ("Notifications", "Real-time alerts for approvals and deadlines."),
    ]

    x_positions = [0.6, 4.55, 8.5]
    y_positions = [1.9, 4.1]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = x_positions[col]
        top = y_positions[row]

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.7), Inches(2.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(3.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT
        p2.font.name = "Inter"
        p2.space_before = Pt(4)
        p2.line_spacing = 1.2

    add_footer(slide, 8)


def add_benefits_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Expected Benefits")

    add_bullet_box(slide, 0.6, 1.9, 5.8, 4.2,
                   ["Centralized record of all OJT students deployed to city offices.",
                    "Faster verification of attendance, reports, and requirements.",
                    "Reduced paperwork and manual follow-ups.",
                    "Better coordination between HR, departments, and schools."],
                   title="For the City Hall")

    add_bullet_box(slide, 6.9, 1.9, 5.8, 4.2,
                   ["Easy discovery of verified OJT openings.",
                    "Direct contact with companies and supervisors.",
                    "Clear view of required hours and approvals.",
                    "All documents and reports in one digital portfolio."],
                   title="For Students")

    add_footer(slide, 9)


def add_summary_slide():
    slide = add_blank_slide()
    add_title_bar(slide, "Summary")

    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.4)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(239, 246, 255)
    highlight.line.color.rgb = PRIMARY

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "InternTrack with OJTConnect is a modern, responsive, and role-based web system designed to solve the OJT management problems faced by colleges and city offices. It helps students find placements, helps supervisors track progress, and helps coordinators and HR oversee the entire program."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Inter"
    p.line_spacing = 1.2

    summary_points = [
        ("Problem", "Manual, fragmented, and slow OJT management process."),
        ("Solution", "Unified web platform with openings, tracking, and role-based workflows."),
        ("Client", "HR Department of the City Hall — the central gateway for OJT deployment."),
        ("Outcome", "Digital, transparent, and scalable internship tracking for the public sector."),
    ]

    for i, (title, desc) in enumerate(summary_points):
        left = 0.6 + (i % 2) * 6.2
        top = 3.7 + (i // 2) * 1.6

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.9), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(5.5), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT
        p2.font.name = "Inter"
        p2.space_before = Pt(2)
        p2.line_spacing = 1.2

    add_footer(slide, 10)


def add_thank_you_slide():
    slide = add_blank_slide()
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(12.1), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "We look forward to partnering with the HR Department of the City Hall to modernize OJT placement and tracking."
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.name = "Inter"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    p3 = tf.add_paragraph()
    p3.text = "Questions & Discussion"
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.font.name = "Inter"
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)

    add_footer(slide, 11)


# Build the presentation
add_cover_slide()
add_title_defense_slide()
add_problem_statement_slide()
add_objectives_slide()
add_scope_slide()
add_why_hr_slide()
add_solution_slide()
add_features_slide()
add_benefits_slide()
add_summary_slide()
add_thank_you_slide()

output_path = "interntrack_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
